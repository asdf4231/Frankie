"""Extract chat attachments into model-ready text or image blocks."""

from __future__ import annotations

import base64
from io import BytesIO
from pathlib import Path

MAX_ATTACHMENT_BYTES = 20 * 1024 * 1024
MAX_EXTRACTED_CHARS = 120_000
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg"}
IMAGE_MIME_TYPES = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg"}


def prepare_attachment(filename: str, data: bytes) -> tuple[str, dict | str]:
    """Return a display name and an Anthropic-compatible content block."""
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        allowed = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"不支持的文件类型，请上传：{allowed}")
    if not data:
        raise ValueError(f"文件为空：{filename}")
    if len(data) > MAX_ATTACHMENT_BYTES:
        raise ValueError(f"文件超过 20MB 上限：{filename}")

    if suffix in IMAGE_MIME_TYPES:
        encoded = base64.b64encode(data).decode("ascii")
        return filename, {
            "type": "image",
            "source": {"type": "base64", "media_type": IMAGE_MIME_TYPES[suffix], "data": encoded},
        }

    if suffix == ".pdf":
        from pypdf import PdfReader

        text = "\n\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(data)).pages)
    elif suffix == ".docx":
        from docx import Document

        document = Document(BytesIO(data))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        for table in document.tables:
            text += "\n" + "\n".join(" | ".join(cell.text for cell in row.cells) for row in table.rows)
    elif suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(BytesIO(data))
        text = "\n\n".join(
            "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            for slide in presentation.slides
        )
    text = text.strip()
    if not text:
        text = "（未提取到文字内容，可能是扫描件或图片型文档）"
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS] + "\n[附件文字已截断]"
    return filename, f"【附件：{filename}】\n{text}"
